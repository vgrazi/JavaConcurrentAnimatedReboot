from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---- palette (matches the app: black bg, green=runnable, blue=keyword, red=terminated) ----
BG      = RGBColor(0x0D, 0x0D, 0x0D)
PANEL   = RGBColor(0x16, 0x16, 0x18)
WHITE   = RGBColor(0xF2, 0xF2, 0xF2)
GREY    = RGBColor(0xB5, 0xB5, 0xBC)
GREEN   = RGBColor(0x2E, 0xD1, 0x6A)
BLUE    = RGBColor(0x4F, 0x9DFF & 0xFF, 0xFF) if False else RGBColor(0x4F, 0x9D, 0xFF)
RED     = RGBColor(0xFF, 0x5A, 0x5A)
YELLOW  = RGBColor(0xF5, 0xD0, 0x42)
MONO    = "Consolas"
SANS    = "Segoe UI"

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW, SH = prs.slide_width, prs.slide_height
blank = prs.slide_layouts[6]


def bg(slide, color=BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def box(slide, x, y, w, h, fill=None, line=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=6, line_spacing=1.06):
    """runs: list of paragraphs; each paragraph is list of (txt, size, color, bold, font, italic)."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (t, sz, col, bold, fnt, *rest) in para:
            it = rest[0] if rest else False
            r = p.add_run(); r.text = t
            r.font.size = Pt(sz); r.font.color.rgb = col
            r.font.bold = bold; r.font.name = fnt; r.font.italic = it
    return tb


def accent_bar(slide, color=GREEN, x=Inches(0.0), y=Inches(0.0), w=Inches(0.14), h=SH):
    b = slide.shapes.add_shape(1, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = color; b.line.fill.background()
    b.shadow.inherit = False
    return b


def kicker(slide, txt, color=GREEN):
    text(slide, Inches(0.6), Inches(0.45), Inches(11), Inches(0.4),
         [[(txt.upper(), 13, color, True, SANS)]])


def title(slide, txt, y=Inches(0.8), size=30, color=WHITE):
    text(slide, Inches(0.6), y, Inches(12.2), Inches(1.1),
         [[(txt, size, color, True, SANS)]])


# ============================================================= TITLE SLIDE
s = prs.slides.add_slide(blank); bg(s)
accent_bar(s, GREEN)
text(s, Inches(0.9), Inches(2.1), Inches(11.6), Inches(1.2),
     [[("Structured Concurrency", 46, WHITE, True, SANS)]])
text(s, Inches(0.9), Inches(3.15), Inches(11.6), Inches(0.7),
     [[("A guided walkthrough of the animation — talking points, button by button",
        20, GREEN, False, SANS)]])
text(s, Inches(0.9), Inches(4.15), Inches(11.6), Inches(1.6),
     [[("Java 25 preview  •  JEP 505  •  java.util.concurrent.StructuredTaskScope", 16, GREY, False, MONO)],
      [("Part of Java Concurrent Animated — Reboot", 14, GREY, False, SANS)]])
# little legend
text(s, Inches(0.9), Inches(6.4), Inches(11.6), Inches(0.6),
     [[("Colour key:  ", 13, GREY, True, SANS),
       ("green = runnable   ", 13, GREEN, True, SANS),
       ("white = waiting   ", 13, WHITE, True, SANS),
       ("blue = keyword   ", 13, BLUE, True, SANS),
       ("red = failed / cancelled", 13, RED, True, SANS)]])


# ============================================================= MENTAL MODEL
s = prs.slides.add_slide(blank); bg(s); accent_bar(s, GREEN)
kicker(s, "The one idea to hold onto")
title(s, "The scope is the box; subtasks live and die inside it")
text(s, Inches(0.6), Inches(1.85), Inches(7.2), Inches(5.2),
     [[("In the animation, the tall grey monolith is the ", 18, WHITE, False, SANS),
       ("StructuredTaskScope", 18, YELLOW, True, MONO),
       (".", 18, WHITE, False, SANS)],
      [("", 8, WHITE, False, SANS)],
      [("• Every subtask you fork is a virtual thread that runs ", 17, WHITE, False, SANS),
       ("inside", 17, GREEN, True, SANS),
       (" that box.", 17, WHITE, False, SANS)],
      [("• The scope owns its children: it cannot be closed until they all finish.", 17, WHITE, False, SANS)],
      [("• If the scope shuts down, the children are ", 17, WHITE, False, SANS),
       ("interrupted and cancelled", 17, RED, True, SANS),
       (" — no leaks, no orphans.", 17, WHITE, False, SANS)],
      [("", 8, WHITE, False, SANS)],
      [("This parent-child confinement is the whole point: concurrency that follows the "
        "shape of your code, so errors and cancellation propagate the way exceptions do.",
        17, GREY, False, SANS, True)]])
# code panel
box(s, Inches(8.05), Inches(1.85), Inches(4.75), Inches(4.3), fill=PANEL, line=RGBColor(0x33,0x33,0x38))
text(s, Inches(8.3), Inches(2.05), Inches(4.3), Inches(4.0),
     [[("try (var scope =", 14, BLUE, True, MONO)],
      [("      StructuredTaskScope.open(J)) {", 14, WHITE, False, MONO)],
      [("  var a = scope.fork(task1);", 14, GREEN, False, MONO)],
      [("  var b = scope.fork(task2);", 14, GREEN, False, MONO)],
      [("  var r = scope.join();", 14, YELLOW, True, MONO)],
      [("  process(r);", 14, WHITE, False, MONO)],
      [("}  // all children done here", 14, GREY, False, MONO)]],
     line_spacing=1.25)


# ============================================================= BUTTON SLIDES
def button_slide(idx, btn, accent, subtitle, points, code=None, note=None):
    s = prs.slides.add_slide(blank); bg(s); accent_bar(s, accent)
    kicker(s, f"Button {idx}", accent)
    # button pill
    pill = box(s, Inches(0.6), Inches(0.85), Inches(6.6), Inches(0.72),
               fill=PANEL, line=accent, line_w=1.5)
    tf = pill.text_frame; tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = btn; r.font.name = MONO; r.font.size = Pt(18)
    r.font.bold = True; r.font.color.rgb = WHITE
    text(s, Inches(0.6), Inches(1.72), Inches(12.1), Inches(0.5),
         [[(subtitle, 16, accent, True, SANS)]])
    top = Inches(2.35)
    body_w = Inches(12.1) if code is None else Inches(7.1)
    runs = []
    for head, sub in points:
        runs.append([("•  ", 16, accent, True, SANS), (head, 16, WHITE, True, SANS)])
        if sub:
            runs.append([("    ", 16, WHITE, False, SANS), (sub, 15, GREY, False, SANS)])
    text(s, Inches(0.6), top, body_w, Inches(4.4), runs, space_after=7, line_spacing=1.08)
    if code:
        box(s, Inches(8.05), top, Inches(4.75), Inches(3.9), fill=PANEL,
            line=RGBColor(0x33,0x33,0x38))
        crun = [[(ln[0], 13.5, ln[1], ln[2], MONO)] for ln in code]
        text(s, Inches(8.3), top + Inches(0.2), Inches(4.3),
             Inches(3.6), crun, line_spacing=1.22)
    if note:
        text(s, Inches(0.6), Inches(6.55), Inches(12.1), Inches(0.7),
             [[("Watch the animation:  ", 14, accent, True, SANS),
               (note, 14, GREY, False, SANS, True)]])
    return s


button_slide(
    1, "open(allSuccessfulOrThrow)", GREEN,
    "Open a scope that demands every subtask succeed",
    [("Creates the scope with the fail-fast \u201call must succeed\u201d joiner.",
      "This is the everyday case: a fan-out where you need every result to proceed."),
     ("join() will return a Stream of the completed subtasks.",
      "You map each Subtask::get to its value \u2014 all guaranteed present."),
     ("The moment ANY subtask throws, the scope cancels the rest.",
      "Fail-fast: no time wasted finishing work whose result you\u2019ll discard."),
     ("Nothing is forked yet \u2014 the box (scope) simply appears, ready.",
      "Owner thread is now parked, waiting for your fork commands.")],
    code=[("StructuredTaskScope.open(", BLUE, True),
          ("  Joiner.<T>allSuccessfulOrThrow())", WHITE, False),
          ("", WHITE, False),
          ("// join() -> Stream<Subtask<T>>", GREY, False),
          ("// one failure -> cancel siblings", GREY, False)],
    note="the grey monolith is drawn \u2014 an empty scope waiting for children.")

button_slide(
    2, "open(anySuccessfulResultOrThrow)", BLUE,
    "Open a scope where the first success wins",
    [("Creates the scope with the \u201crace\u201d joiner.",
      "Perfect for redundancy: ask three mirrors, take whoever answers first."),
     ("join() returns the single winning value directly \u2014 not a stream.",
      "As soon as one subtask succeeds, join() unblocks with that result."),
     ("The losing siblings are cancelled automatically.",
      "You never wait for the slower replicas; they\u2019re interrupted."),
     ("Only fails if EVERY subtask fails.",
      "Then join() throws FailedException \u2014 there was no success to return.")],
    code=[("StructuredTaskScope.open(", BLUE, True),
          ("  Joiner.<T>anySuccessfulResultOrThrow())", WHITE, False),
          ("", WHITE, False),
          ("// join() -> T (the winner)", GREY, False),
          ("// first success cancels the rest", GREY, False)],
    note="same empty scope \u2014 but now wired for first-one-wins semantics.")

button_slide(
    3, "fork(subtask)", GREEN,
    "Add a child task that runs concurrently inside the scope",
    [("scope.fork(...) starts a new virtual thread for the subtask.",
      "In the animation a rounded runner enters the box and parks (white = in-flight)."),
     ("Fork as many as you like \u2014 the scope grows to hold them all.",
      "Each child is cheap: virtual threads scale to thousands."),
     ("Forking is confined to the owner thread, before join().",
      "You cannot fork after join() \u2014 the animation blocks that, just like the API."),
     ("A returned Subtask handle lets you read its value after join().",
      "But you never touch it until join() says the results are ready.")],
    code=[("Subtask<String> s1 =", WHITE, False),
          ("    scope.fork(() -> fetchA());", GREEN, False),
          ("Subtask<String> s2 =", WHITE, False),
          ("    scope.fork(() -> fetchB());", GREEN, False),
          ("", WHITE, False),
          ("// children now live in the scope", GREY, False)],
    note="each click adds a subtask sprite parked inside the monolith.")

button_slide(
    4, "Succeed a subtask", GREEN,
    "Let one in-flight child complete normally",
    [("Resolves the oldest pending subtask with a value.",
      "This is a teaching control \u2014 it stands in for the child\u2019s work finishing."),
     ("In allSuccessful mode: the scope keeps waiting for the others.",
      "join() only returns once every child has succeeded."),
     ("In anySuccessful mode: this is the winner \u2014 join() returns it.",
      "and the remaining siblings are immediately cancelled."),
     ("The succeeding sprite exits the scope to the right (done).",
      "Green \u2192 terminated: it carried a result out of the box.")],
    note="the resolved runner leaves the box to the right, carrying its value.")

button_slide(
    5, "Fail a subtask", RED,
    "Make one child throw \u2014 and watch cancellation propagate",
    [("The oldest pending subtask throws an exception.",
      "This is where structured concurrency truly earns its keep."),
     ("allSuccessful mode: fail-fast \u2014 all siblings are cancelled at once.",
      "Their virtual threads are interrupted; the sprites retreat to the left (red)."),
     ("anySuccessful mode: a single failure is tolerated.",
      "The scope keeps racing; it only throws if ALL children fail."),
     ("The failure is remembered and surfaced by join() as FailedException.",
      "No exception is ever silently lost, unlike a detached future.")],
    code=[("// inside a subtask", GREY, False),
          ("throw new IOException(...);", RED, True),
          ("", WHITE, False),
          ("// allSuccessful:", GREY, False),
          ("//   -> siblings interrupted", RED, False),
          ("//   -> join() throws Failed", RED, False)],
    note="siblings retreat to the left the instant one child fails (fail-fast).")

button_slide(
    6, "join()", YELLOW,
    "Wait for the scope to reach its conclusion \u2014 the join point",
    [("Blocks the owner thread until the joiner\u2019s condition is met.",
      "all-succeeded, first-success, or a failure \u2014 whichever comes first."),
     ("On success: returns results and the NEXT STEP runs.",
      "The animation sends a downstream \u201cnext step\u201d thread through the scope."),
     ("On failure: throws FailedException; the next step is skipped.",
      "Control jumps to your catch block \u2014 exactly like ordinary sequential code."),
     ("join() re-establishes a single thread of control after the fan-out.",
      "Everything before it happened concurrently; everything after is sequential again.")],
    code=[("try {", BLUE, True),
          ("  var r = scope.join();", YELLOW, True),
          ("  process(r);   // next step", GREEN, False),
          ("} catch (FailedException e) {", BLUE, True),
          ("  handle(e.getCause());", RED, False),
          ("}", BLUE, True)],
    note="green \u201cnext step\u201d runs on success; on failure the catch line lights up.")

button_slide(
    7, "Reset", GREY,
    "Tear the scope down and start over",
    [("Interrupts the owner thread and closes any open scope.",
      "try-with-resources guarantees the scope is closed and children are cancelled."),
     ("Clears all subtasks and returns the box to its empty state.",
      "A clean slate to demo the other joiner or another fork pattern."),
     ("Demonstrates deterministic cleanup \u2014 the core promise.",
      "When the block exits, there are provably no threads still running.")],
    note="the monolith empties and shrinks back \u2014 nothing left running.")


# ============================================================= SC vs CF (why)
s = prs.slides.add_slide(blank); bg(s); accent_bar(s, GREEN)
kicker(s, "Why it matters")
title(s, "Benefits over CompletableFuture")
benefits = [
    ("Readable, sequential shape", "Fork/join lives in one try block. Control flow reads top-to-bottom instead of a chain of thenCompose / thenCombine callbacks."),
    ("No leaked threads", "The scope can\u2019t close until every child finishes or is cancelled. A stray CompletableFuture can outlive its caller and run forever."),
    ("Automatic cancellation", "One failure cancels the siblings for free. With CompletableFuture you must wire up cancellation and propagation by hand."),
    ("Errors can\u2019t vanish", "Every subtask exception is collected and surfaced by join(). A CompletableFuture whose result is never read can swallow its exception."),
    ("Real stack traces & observability", "Children are ordinary threads with a parent link, so thread dumps show the tree. CF work hops around a shared pool, hiding structure."),
    ("Cheap fan-out via virtual threads", "Blocking calls in subtasks are fine \u2014 no need to make everything non-blocking just to avoid starving a pool."),
]
col_w = Inches(6.0)
for i, (h, sub) in enumerate(benefits):
    col = i % 2; row = i // 2
    x = Inches(0.6) + (col * Inches(6.35))
    y = Inches(1.85) + (row * Inches(1.62))
    box(s, x, y, col_w, Inches(1.42), fill=PANEL, line=RGBColor(0x2C,0x2C,0x30))
    text(s, x + Inches(0.22), y + Inches(0.12), col_w - Inches(0.44), Inches(1.2),
         [[(h, 16, GREEN, True, SANS)],
          [(sub, 12.5, GREY, False, SANS)]], line_spacing=1.05, space_after=3)


# ============================================================= When CF is still better
s = prs.slides.add_slide(blank); bg(s); accent_bar(s, BLUE)
kicker(s, "Use the right tool", BLUE)
title(s, "When CompletableFuture is still the better choice")
cases = [
    ("Fire-and-forget / long-lived work", "Tasks that must outlive the current method \u2014 background refresh, caches, work that isn\u2019t tied to one request\u2019s lifetime. A scope insists on joining before it exits."),
    ("Event-driven pipelines & composition", "Rich chaining with thenApply / thenCompose / thenCombine, or reacting whenComplete. CF is a composable value; a scope is a lifetime, not a pipeline stage."),
    ("A result others subscribe to later", "You want to hand back a CompletableFuture that many callers attach callbacks to over time. Structured scopes are consumed at the join point, not passed around."),
    ("Callback / non-blocking APIs", "Bridging an async, callback-based library via completeExceptionally / complete fits CF naturally."),
    ("Pre-Java 21, or no preview flag", "StructuredTaskScope is a preview API. CompletableFuture has been stable since Java 8 and needs no --enable-preview."),
]
y = Inches(1.85)
for h, sub in cases:
    box(s, Inches(0.6), y, Inches(12.15), Inches(0.98), fill=PANEL, line=RGBColor(0x2C,0x2C,0x30))
    text(s, Inches(0.85), y + Inches(0.1), Inches(11.7), Inches(0.8),
         [[(h, 15.5, BLUE, True, SANS),
           ("   \u2014 " + sub, 13, GREY, False, SANS)]], line_spacing=1.02)
    y = y + Inches(1.06)


# ============================================================= RULE OF THUMB
s = prs.slides.add_slide(blank); bg(s); accent_bar(s, GREEN)
kicker(s, "Takeaway")
title(s, "Rule of thumb")
box(s, Inches(0.6), Inches(2.0), Inches(12.15), Inches(1.5), fill=PANEL, line=GREEN, line_w=1.5)
text(s, Inches(0.95), Inches(2.15), Inches(11.5), Inches(1.2),
     [[("Concurrent sub-work for one task, that must all finish before you continue?  ",
        18, WHITE, False, SANS),
       ("Structured Concurrency.", 18, GREEN, True, SANS)]],
     anchor=MSO_ANCHOR.MIDDLE)
box(s, Inches(0.6), Inches(3.7), Inches(12.15), Inches(1.5), fill=PANEL, line=BLUE, line_w=1.5)
text(s, Inches(0.95), Inches(3.85), Inches(11.5), Inches(1.2),
     [[("A composable, long-lived, or subscribed-to async value?  ",
        18, WHITE, False, SANS),
       ("CompletableFuture.", 18, BLUE, True, SANS)]],
     anchor=MSO_ANCHOR.MIDDLE)
text(s, Inches(0.6), Inches(5.6), Inches(12.15), Inches(1.2),
     [[("They\u2019re complementary: reach for structured concurrency to orchestrate a fan-out "
        "within a request, and keep CompletableFuture for reactive composition and work that "
        "lives beyond the call.", 16, GREY, False, SANS, True)]])

out = "/app/JavaConcurrentAnimatedReboot/docs/StructuredConcurrency-TalkingPoints.pptx"
import os
os.makedirs(os.path.dirname(out), exist_ok=True)
prs.save(out)
print("saved", out, "slides:", len(prs.slides._sldIdLst))
